import os
from groq import Groq
from pdfminer.high_level import extract_text
from pptx import Presentation
from PIL import Image
import pytesseract
import cv2

# Keep it highly secure (DO NOT SHARE IT WITH ANYONE!)
client = Groq(
    api_key="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",
)

AI_tutor_prompt = "You are an AI tutor that is focused on helping students by explaining them their notes when asked. You are friendly and use simple language. You will do whatever the student asks you to do. Make sure to answer as correctly as possible and answer only what the student has asked. "
AI_evaluator_prompt = "You are an AI that generates quizzes based on the notes of the student. You generate at least 5 very creative questions along with their answers that tests each and every part of student's understanding of the notes. You do whatever you are asked to do. Here are some rules you MUST FOLLOW: 1) The output you generate must be in JSON format. 2) There should be nothing except JSON format. 3) The name of main key should be 'test'. 4) The value of 'test' should be a dictionary whose key will be the question and value will be the list of answers. 5) In the list of answers there should be 4 random answers with only 1 of them being correct. 6) Make a subkey called 'answers' whose value should be list of correct answers in order of the questions. 7) DO NOT SHOW ANSWERS IN THE OPTIONS YOU PROVIDE! 8) Generate unique and different questions. 9) Do not disobey the above specified rules."
AI_written_evaluator_prompt = "You are an AI that generates 5 questions (long answer type) according to notes of the student and using internet. Here are some rules you must follow 1) The output you generate must be in JSON format. 2) There should be nothing except JSON format. 3) The name of main key should be 'test'. 4) The value of 'test' should be a dictionary whose key will be the question and value will be the list of answers which you think would be appropriate. There should be only 1 answer in each list. 5) Make a subkey called 'answers' whose value should be list of answers in order of the questions. 6) Try adding at least 1 question which is NOT from the notes. 7) Do not disobey the above specified rules."
AI_quiz_evaluator_prompt = "You are an AI that reads student's MCQ quiz data and their notes and then suggests them which part of their notes should they focus more on. You also tell them what they are doing wrong and how they can improve themselves. You are very friendly and helpful. The output you generate should be in beautiful html format for easy understanding of student."
AI_compare_answers_prompt = "You are an AI that compares student's answer with the answer provided and tell if the student's answer is correct or not. Here are some rules you should follow 1) If the answer is correct you should respond 'c' 2) if the answer is wrong you should respond 'w' 3) You should not respond anything else. 4) If the student's answer is empty give him automatically 'w'."
AI_researcher_prompt = "You are an AI professor that generates 1 very unique and outstanding topic from the notes that the student provides for him to do research on. The topic you provide must make the student understand the notes very thoroughly and should help him. Give student some hints and resources on how he can proceed with the research. Make sure the topic of research is easy and not complex and anyone can do it within very less time. The output you generate should be in beautiful html format for easy understanding of student."
AI_project_prompt = "You are an AI that generates a single and unique project topic from the notes of the students that covers the main understanding of the notes. Here are the examples for your reference. If a student uploads notes about programming, the output should be something similar to 'build a banking application using the programming language and upload the report of your project.' If the student uploads a note on hardware related subject, the output should be something similar to 'Build a robot that can utilize cameras to look around itself'. Here are some rules you must follow 1) The output should only be the details about the project and some hints on how the student can start with it. 2)The output you generate should be in beautiful html format for easy understanding of student."
AI_project_evaluator_prompt = "You are an AI that evaluates the project report submitted by student. You use 3 main metrics to evaluate the project 1) Closeness of student's project to the question 2) Creativity of the student 3) Understanding of the student"
AI_notes_detail_extractor_prompt = "You are an AI and you have to answer the following questions in 'yes' or 'no' only. Do you think the following notes 1) require sutdents to make research paper on it? 2) require practical projects? 3) Require group practical projects? 4) what is the quality score of notes on scale of 100. Here are some rules you must follow 1) The answer should be in the format of a list containing yes or no only and nothing else. 2) The output should be a list and nothing other than a list containing 4 elements. 4) Do not give practical projects for theoretical subjects. 5) Be very strict while giving quality score. 6) Do not disobey the above specified rules."
AI_important_questions_prompt = "You are an AI that generates as many questions and answers as you can based on the notes that the student provides. Your job is to find out most important questions from the notes that you believe will appear in exam or are very crucial for knowledge of the student along with their answers. Here are some rules you must follow: 1) The output should be an html code that only has paragraphs and nothing else. 2) Generate 10 questions along with their answers."
AI_research_evaluator_prompt = "You are an AI that evaluates the research paper submitted by student. You use 3 main metrics to evaluate the research paper 1) Closeness of student's research to the question 2) Creativity of the student 3) Understanding of the student. Here are some rules you must follow 1) The output you generate must be beautiful and in html format. 2) These are the elements that should be present in your output 'score' out of 100, be very strict while evaluating the student, list of weak points of the student telling him in which part of the report he is weak in and what should he do to overcome the weakness, lastly your overall opinion on how the work submitted by student was. 3) Do not disobey the above specified rules. "
AI_project_evaluator_prompt = "You are an AI that evaluates the projects submitted by student. You use 3 main metrics to evaluate the projects 1) Closeness of student's research to the question 2) Creativity of the student 3) Understanding of the student. Here are some rules you must follow 1) The output you generate must be beautiful, eye catching, colorful and in html format. 2) These are the elements that should be present in your output 'score' out of 100, be very strict while evaluating the student, list of weak points of the student telling him in which part of the report he is weak in and what should he do to overcome the weakness, lastly your overall opinion on how the work submitted by student was. 3) Do not disobey the above specified rules. "

def respond(task):
    info = {"role": "user",
            "content": task}
    chat_completion = client.chat.completions.create(
        messages=[info],
        model="mixtral-8x7b-32768",
    )
    return chat_completion.choices[0].message.content

def ask_tutor(content):
    AI = "You are an AI tutor that is focused on helping students by explaining them their notes or making them short quizzes based on their notes. You are friendly and use simple language. You will do whatever the user asks you to do."
    info = {"role": "user", 
            "content": AI + " From here on all of this will be user input: " + content}
    chat_completion = client.chat.completions.create(
        messages=[info],
        model="mixtral-8x7b-32768",
    )
    return chat_completion.choices

def make_schedule(content):
    AI = "You are an AI scheduler that is focused on helping students by making schedules based on their subjects so that he does not face any pressure or has to attend too many quizzez in one day. You are friendly and use simple language. You will do whatever the user asks you to do."
    info = {"role": "user", 
            "content": AI + " From here on all of this will be user input: " + content}
    chat_completion = client.chat.completions.create(
        messages=[info],
        model="mixtral-8x7b-32768",
    )
    return chat_completion.choices[0].message.content

def evaluate_student(content):
    AI = "You are an AI evaluater that evaluates student based on how they perform. Your job is to make quizzez or exams related to student's notes that they will solve and you will have to evaluate based on that. You are very creative. Come up with new ideas that are unique to evaluate students and then ask the user these questions that you came up with."
    AI2 = "You are an AI evaluater that evaluates student based on how they perform. Your job is to make quizzez or exams related to student's notes that they will solve and you will have to evaluate based on that. You are very creative. Come up with new ideas that are unique to evaluate student."
    info = {"role": "user", 
            "content": AI2 + " From here on all of this will be user input: " + content}
    chat_completion = client.chat.completions.create(
        messages=[info],
        model="mixtral-8x7b-32768",
    )
    return chat_completion.choices[0].message.content

def clean_text(text):
    # Remove newline characters
    text = text.replace('\n', ' ')
    # Remove extra spaces
    text = ' '.join(text.split())
    return text

def extract_text_from_pdf(pdf_file):
    text = extract_text(pdf_file)
    return clean_text(text)

def extract_text_from_ppt(ppt_file):
    # Load the PowerPoint file
    prs = Presentation(ppt_file)

    # Initialize an empty string to store all text
    text = ""

    # Loop through each slide in the presentation
    for slide in prs.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return clean_text(text)

def extract_text_from_txt(txt_file):
    with open(txt_file, 'r', encoding='utf-8') as file:
        text = file.read()
    return clean_text(text)

def extract_text_from_image(image_file):
    # Open the image file
    #img = Image.open(image_file)

    img = cv2.imread(image_file)  # Load the image
    img = cv2.resize(img, (0, 0), fx=0.25, fy=0.25)
    img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)  # convert to grey
    img = cv2.adaptiveThreshold(img, 255, cv2.ADAPTIVE_THRESH_MEAN_C, cv2.THRESH_BINARY, 3, 15)

    # Use pytesseract to extract text
    text = pytesseract.image_to_string(img)
    # Clean the extracted text
    return clean_text(text)

def extract_text_from_file(file_path):
    # Get the file extension
    _, file_extension = os.path.splitext(file_path)
    
    # Determine the extraction type based on the file extension
    if file_extension in ['.pdf']:
        return extract_text_from_pdf(file_path)
    elif file_extension in ['.ppt', '.pptx']:
        return extract_text_from_ppt(file_path)
    elif file_extension in ['.txt']:
        return extract_text_from_txt(file_path)
    elif file_extension in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
        return extract_text_from_image(file_path)
    else:
        return 'Unsupported file type'



